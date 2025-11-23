# FINSIGHT CORE V3.1 - ENTERPRISE EDITION (CLEAN & POWERFUL)
# Includes: PPTX Export, Deep Dive Drill-Downs, Scenario Planning, Expanded KPIs.
# Fixes: Matplotlib dependency for styling, removed emojis from tabs.

import streamlit as st
import pandas as pd
import numpy as np
import requests
import json
import plotly.graph_objects as go
import plotly.express as px
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# --- CONFIGURATION & SETUP ---
st.set_page_config(layout="wide", page_title="FinSight Core | Enterprise")

# --- CUSTOM CSS ---
st.markdown("""
<style>
    /* 1. Main Gradient Background */
    .stApp {
        background: linear-gradient(135deg, #1e293b 0%, #0f172a 100%);
        background-attachment: fixed;
    }
    
    /* 2. Glass Card Style */
    .glass-card {
        background: rgba(255, 255, 255, 0.05);
        box-shadow: 0 8px 32px 0 rgba(0, 0, 0, 0.37);
        backdrop-filter: blur(8px);
        -webkit-backdrop-filter: blur(8px);
        border-radius: 12px;
        border: 1px solid rgba(255, 255, 255, 0.1);
        padding: 20px;
        margin-bottom: 20px;
        color: white;
    }

    /* 3. Text Styling */
    h1, h2, h3, h4, p, div, span, label, .stMarkdown {
        color: #e2e8f0 !important;
        font-family: 'Inter', sans-serif;
    }
    
    /* 4. Tab Styling - CLEAN (No Emojis) */
    .stTabs [data-baseweb="tab-list"] { gap: 8px; }
    .stTabs [data-baseweb="tab"] {
        height: 45px;
        background-color: rgba(255, 255, 255, 0.05);
        border-radius: 8px 8px 0 0;
        color: #94a3b8;
        font-weight: 600;
        border: none;
    }
    .stTabs [aria-selected="true"] {
        background-color: rgba(255, 255, 255, 0.15);
        color: #ffffff;
        border-bottom: 2px solid #38bdf8;
    }
    
    /* 5. Metric Styling */
    .metric-label { font-size: 0.75rem; color: #94a3b8; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 5px; }
    .metric-value { font-size: 1.6rem; font-weight: 700; color: #ffffff; }
    .metric-delta-pos { color: #4ade80; font-size: 0.85rem; font-weight: 600; }
    .metric-delta-neg { color: #f87171; font-size: 0.85rem; font-weight: 600; }
    
    /* 6. Table Styling Fix */
    .dataframe { font-size: 0.8rem !important; color: white !important; }
</style>
""", unsafe_allow_html=True)

# Gemini API Configuration
GEMINI_MODEL = 'gemini-2.5-flash-preview-09-2025'
API_URL = f"https://generativelanguage.googleapis.com/v1beta/models/{GEMINI_MODEL}:generateContent"

# --- EXPANDED MOCK DATA ---
def generate_mock_data():
    # Base Monthly Data
    months = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'Jul', 'Aug', 'Sep', 'Oct', 'Nov', 'Dec']
    
    data = []
    for i, m in enumerate(months):
        is_actual = i < 6
        # Base values with some randomness
        rev = 120000 + (i * 5000) + np.random.randint(-5000, 5000)
        direct = rev * 0.4 + np.random.randint(-2000, 2000)
        indirect = 40000 + (i * 1000) + np.random.randint(-1000, 1000)
        cash_flow = (rev - direct - indirect) * 0.9 # Lag
        
        data.append({
            'month': m,
            'type': 'Actual' if is_actual else 'Forecast',
            'revenue': rev,
            'direct_cost': direct,
            'indirect_cost': indirect,
            'cash_flow': cash_flow,
            'headcount': 120 + (i * 2),
            'budget_revenue': 120000 + (i * 4000),
            'budget_direct': (120000 + (i * 4000)) * 0.38,
            'budget_indirect': 42000
        })
    return pd.DataFrame(data)

ANNUAL_DATA = generate_mock_data()
ANNUAL_DATA['gross_margin'] = ANNUAL_DATA['revenue'] - ANNUAL_DATA['direct_cost']
ANNUAL_DATA['opex'] = ANNUAL_DATA['indirect_cost']
ANNUAL_DATA['operating_profit'] = ANNUAL_DATA['gross_margin'] - ANNUAL_DATA['opex']
ANNUAL_DATA['operating_margin_pct'] = ANNUAL_DATA['operating_profit'] / ANNUAL_DATA['revenue']

# Drill Down Data (Mock) - EXPANDED
DRILL_DATA = pd.DataFrame([
    {'Category': 'Revenue', 'Sub-Category': 'Product A', 'Cost Centre': 'Sales-US', 'Account': '4001-Sales', 'Transaction': 'Inv-1001', 'Actual': 450000, 'Budget': 420000},
    {'Category': 'Revenue', 'Sub-Category': 'Product B', 'Cost Centre': 'Sales-EU', 'Account': '4001-Sales', 'Transaction': 'Inv-1002', 'Actual': 320000, 'Budget': 350000},
    {'Category': 'COGS', 'Sub-Category': 'Materials', 'Cost Centre': 'Plant-1', 'Account': '5001-Mat', 'Transaction': 'PO-5501', 'Actual': 200000, 'Budget': 190000},
    {'Category': 'COGS', 'Sub-Category': 'Labor', 'Cost Centre': 'Plant-1', 'Account': '5005-Wages', 'Transaction': 'Payroll-06', 'Actual': 150000, 'Budget': 145000},
    {'Category': 'OPEX', 'Sub-Category': 'Marketing', 'Cost Centre': 'Mktg-Global', 'Account': '6100-Agency', 'Transaction': 'Inv-9901', 'Actual': 80000, 'Budget': 70000},
    {'Category': 'OPEX', 'Sub-Category': 'R&D', 'Cost Centre': 'Eng-HQ', 'Account': '6200-Cloud', 'Transaction': 'AWS-Bill', 'Actual': 120000, 'Budget': 120000},
    {'Category': 'OPEX', 'Sub-Category': 'G&A', 'Cost Centre': 'Corp-HQ', 'Account': '6300-Legal', 'Transaction': 'Inv-Legal', 'Actual': 50000, 'Budget': 45000},
])

# --- HELPER FUNCTIONS ---

def format_k(val):
    return f"£{val/1000:.1f}k"

def format_pct(val):
    return f"{val*100:.1f}%"

def render_glass_metric(label, value, delta, is_good=True):
    delta_class = "metric-delta-pos" if is_good else "metric-delta-neg"
    arrow = "▲" if is_good else "▼"
    st.markdown(f"""
    <div class="glass-card" style="padding: 15px; height: 100%;">
        <div class="metric-label">{label}</div>
        <div class="metric-value">{value}</div>
        <div class="{delta_class}">{arrow} {delta} <span style="color: #64748b; font-size: 0.7rem; font-weight: 400;">vs Budget</span></div>
    </div>
    """, unsafe_allow_html=True)

def apply_glass_style(fig):
    fig.update_layout(
        paper_bgcolor='rgba(255, 255, 255, 0.0)',
        plot_bgcolor='rgba(0,0,0,0)',
        font=dict(color='#e2e8f0'),
        margin=dict(l=20, r=20, t=30, b=20),
        xaxis=dict(showgrid=False, color='#94a3b8'),
        yaxis=dict(showgrid=True, gridcolor='rgba(255,255,255,0.1)', color='#94a3b8')
    )
    return fig

# --- PPTX GENERATOR ---
def create_presentation(metrics, ai_text):
    prs = Presentation()
    
    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "FinSight Executive Summary"
    subtitle.text = "Automated Board Pack | Period: June 2024"
    
    # Slide 2: Key Financials Table
    slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank
    title = slide.shapes.title
    title.text = "Financial Key Performance Indicators"
    
    # Add Table
    rows, cols = 7, 3
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Headers
    headers = ['Metric', 'Value', 'Variance']
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
    
    # Data 
    data = [
        ['Revenue (YTD)', metrics['rev_val'], metrics['rev_delta']],
        ['Gross Margin', metrics['gm_val'], metrics['gm_delta']],
        ['Operating Profit', metrics['op_val'], metrics['op_delta']],
        ['Operating Margin', metrics['op_margin_val'], metrics['op_margin_delta']],
        ['OPEX', metrics['opex_val'], metrics['opex_delta']],
        ['Cash Flow', metrics['cf_val'], metrics['cf_delta']]
    ]
    
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, val in enumerate(row_data):
            table.cell(row_idx, col_idx).text = str(val)

    # Slide 3: AI Commentary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "AI-Driven Strategic Commentary"
    body = slide.placeholders[1]
    # Basic cleanup for PPTX (removes markdown bolding)
    clean_text = ai_text.replace('**', '').replace('###', '').replace('####', '')
    body.text = clean_text
    
    binary_output = BytesIO()
    prs.save(binary_output)
    binary_output.seek(0)
    return binary_output

# --- AI NARRATIVE ---
def generate_ai_narrative():
    try:
        api_key = st.secrets["GEMINI_API_KEY"]
    except KeyError:
        return "⚠️ API Key Missing in Streamlit Secrets."

    # Context with more data
    ytd = ANNUAL_DATA[ANNUAL_DATA['type'] == 'Actual'].sum()
    budget_rev = ytd['budget_revenue']
    
    context = f"""
    YTD Revenue: {format_k(ytd['revenue'])} (Budget: {format_k(budget_rev)}).
    YTD Op Profit: {format_k(ytd['operating_profit'])}.
    Cash Flow: {format_k(ytd['cash_flow'])}.
    Forecast Accuracy: 94.2%.
    Headcount: {int(ANNUAL_DATA.iloc[5]['headcount'])}.
    """
    
    prompt = f"""You are a Group CFO. Write a strict, board-level executive summary based on: {context}.
    Structure:
    1. **Executive Headline**: One punchy sentence summarizing the month.
    2. **Performance Drivers**: Bullet points on Revenue, Margin, and OPEX.
    3. **Risk Radar**: Specific cash or headcount risks.
    4. **Strategic Recommendation**: One forward-looking action.
    Keep it under 200 words. Use Markdown.
    """

    payload = {'contents': [{'parts': [{'text': prompt}]}]}
    
    try:
        response = requests.post(f"{API_URL}?key={api_key}", headers={'Content-Type': 'application/json'}, data=json.dumps(payload))
        if response.status_code == 200:
            return response.json()['candidates'][0]['content']['parts'][0]['text']
        else:
            return f"AI Error: {response.status_code}"
    except Exception as e:
        return f"AI Service Unavailable: {str(e)}"

# --- VIEWS ---

def dashboard_view():
    
    # 1. EXECUTIVE SUMMARY STRIP (Expanded KPIs)
    # Calc YTD
    ytd_act = ANNUAL_DATA[ANNUAL_DATA['type'] == 'Actual'].sum()
    ytd_bud = ANNUAL_DATA[ANNUAL_DATA['type'] == 'Actual'][['budget_revenue', 'budget_direct', 'budget_indirect']].sum()
    
    # Variances
    rev_var = (ytd_act['revenue'] - ytd_bud['budget_revenue']) / ytd_bud['budget_revenue']
    
    gm_act = ytd_act['revenue'] - ytd_act['direct_cost']
    gm_bud = ytd_bud['budget_revenue'] - ytd_bud['budget_direct']
    gm_var = (gm_act - gm_bud) / gm_bud
    
    op_act = ytd_act['operating_profit']
    op_bud = gm_bud - ytd_bud['budget_indirect']
    op_var = (op_act - op_bud) / op_bud
    
    opex_act = ytd_act['opex']
    opex_bud = ytd_bud['budget_indirect']
    opex_var = (opex_act - opex_bud) / opex_bud
    
    cf_var = 0.054 # Mock
    
    # Store for PPTX
    metrics_for_ppt = {
        'rev_val': format_k(ytd_act['revenue']), 'rev_delta': format_pct(rev_var),
        'gm_val': format_k(gm_act), 'gm_delta': format_pct(gm_var),
        'op_val': format_k(op_act), 'op_delta': format_pct(op_var),
        'op_margin_val': format_pct(ytd_act['operating_margin_pct']/6), 'op_margin_delta': format_pct(op_var),
        'opex_val': format_k(opex_act), 'opex_delta': format_pct(opex_var),
        'cf_val': format_k(ytd_act['cash_flow']), 'cf_delta': format_pct(cf_var)
    }

    st.markdown("### 🚀 Enterprise Performance")
    c1, c2, c3, c4, c5, c6 = st.columns(6)
    with c1: render_glass_metric("Revenue", format_k(ytd_act['revenue']), format_pct(rev_var), rev_var > 0)
    with c2: render_glass_metric("Gross Margin", format_k(gm_act), format_pct(gm_var), gm_var > 0)
    with c3: render_glass_metric("Op Margin", format_pct(ytd_act['operating_margin_pct']/6), format_pct(op_var), op_var > 0) 
    with c4: render_glass_metric("OPEX", format_k(opex_act), format_pct(opex_var), opex_var < 0) # Lower opex is good
    with c5: render_glass_metric("Cash Flow", format_k(ytd_act['cash_flow']), format_pct(cf_var), True)
    with c6: render_glass_metric("Accuracy", "94.2%", "-1.1%", False)

    # 2. SCENARIO SIMULATOR (Interactive)
    with st.expander("🎛️ Scenario Simulation & Sensitivity Analysis", expanded=False):
        sc_col1, sc_col2 = st.columns([1, 3])
        with sc_col1:
            st.markdown("#### Assumptions")
            rev_growth = st.slider("Revenue Growth Impact", -20, 20, 0, format="%d%%")
            cogs_inf = st.slider("COGS Inflation", 0, 10, 0, format="%d%%")
        with sc_col2:
            # Apply scenario
            scenario_data = ANNUAL_DATA.copy()
            scenario_data['revenue'] = scenario_data['revenue'] * (1 + rev_growth/100)
            scenario_data['direct_cost'] = scenario_data['direct_cost'] * (1 + cogs_inf/100)
            scenario_data['gross_margin'] = scenario_data['revenue'] - scenario_data['direct_cost']
            
            fig_sc = go.Figure()
            fig_sc.add_trace(go.Bar(x=scenario_data['month'], y=scenario_data['gross_margin'], name='Projected GM', marker_color='#38bdf8'))
            fig_sc.add_trace(go.Scatter(x=ANNUAL_DATA['month'], y=ANNUAL_DATA['gross_margin'], name='Baseline GM', line=dict(color='white', dash='dot')))
            fig_sc = apply_glass_style(fig_sc)
            fig_sc.update_layout(title="Scenario Impact: Gross Margin", height=300)
            st.plotly_chart(fig_sc, use_container_width=True)

    # 3. MAIN CHARTS (Forecast & Waterfall)
    col_main, col_side = st.columns([2, 1])

    with col_main:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        # Combined Chart
        fig = go.Figure()
        fig.add_trace(go.Bar(x=ANNUAL_DATA['month'], y=ANNUAL_DATA['revenue'], name='Revenue', marker_color='#38bdf8', opacity=0.6))
        fig.add_trace(go.Scatter(x=ANNUAL_DATA['month'], y=ANNUAL_DATA['operating_profit'], name='Op Profit', line=dict(color='#4ade80', width=3), yaxis='y2'))
        
        fig.update_layout(
            title="Revenue & Profitability Trend (YTD + Forecast)",
            yaxis=dict(title="Revenue", showgrid=False),
            yaxis2=dict(title="Profit", overlaying='y', side='right', showgrid=False),
            legend=dict(orientation="h", y=1.1),
            height=350
        )
        fig = apply_glass_style(fig)
        st.plotly_chart(fig, use_container_width=True)
        st.markdown('</div>', unsafe_allow_html=True)

    with col_side:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        # Profit Bridge
        fig_bridge = go.Figure(go.Waterfall(
            name = "20", orientation = "v",
            measure = ["relative", "relative", "relative", "total"],
            x = ["Budget Profit", "Vol Variance", "Cost Var", "Actual Profit"],
            textposition = "outside",
            text = ["+10k", "+5k", "-2k", "£68k"], # Mock calc for demo
            y = [55000, 15000, -2000, 0],
            connector = {"line":{"color":"white"}},
            decreasing = {"marker":{"color":"#f87171"}},
            increasing = {"marker":{"color":"#4ade80"}},
            totals = {"marker":{"color":"#94a3b8"}}
        ))
        fig_bridge.update_layout(title="Profitability Bridge (June)", height=350)
        fig_bridge = apply_glass_style(fig_bridge)
        st.plotly_chart(fig_bridge, use_container_width=True)
        st.markdown('</div>', unsafe_allow_html=True)

    # 4. DEEP DIVE DRILL-DOWNS (Enhanced)
    st.markdown("### 🔍 Multi-Layer Drill Down")
    col_d1, col_d2 = st.columns([1, 2])
    
    with col_d1:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        level = st.selectbox("Hierarchy Level", ["Category", "Sub-Category", "Cost Centre", "Account", "Transaction"])
        
        # Dynamic filtering based on level (Simulated)
        if level == "Category":
            options = DRILL_DATA['Category'].unique()
        elif level == "Sub-Category":
            options = DRILL_DATA['Sub-Category'].unique()
        else:
            options = DRILL_DATA['Category'].unique() # Fallback for demo
            
        selected_filter = st.selectbox(f"Select {level}", options)
        st.markdown('</div>', unsafe_allow_html=True)
        
    with col_d2:
        # Simple filter logic for demo purposes
        if level == "Category":
            filtered_df = DRILL_DATA[DRILL_DATA['Category'] == selected_filter]
        elif level == "Sub-Category":
            filtered_df = DRILL_DATA[DRILL_DATA['Sub-Category'] == selected_filter]
        else:
            filtered_df = DRILL_DATA # Show all for other levels in demo
            
        filtered_df = filtered_df.copy()
        filtered_df['Variance'] = filtered_df['Actual'] - filtered_df['Budget']
        
        # Styled Table
        st.dataframe(
            filtered_df.style.format("{:,.0f}", subset=['Actual', 'Budget', 'Variance'])
            .background_gradient(cmap='RdYlGn_r', subset=['Actual']) # Just to show gradient working
            , use_container_width=True
        )

    # 5. AI & EXPORT (The "CFO" Close)
    st.markdown("### 🤖 Strategic Intelligence & Reporting")
    
    col_ai, col_export = st.columns([3, 1])
    
    with col_ai:
        if 'ai_analysis' not in st.session_state:
            st.session_state['ai_analysis'] = "Click 'Generate' to run the Commentary Engine."
            
        if st.button("⚡ Run Commentary Engine"):
            with st.spinner("Aggregating variances and writing board memo..."):
                st.session_state['ai_analysis'] = generate_ai_narrative()
        
        st.markdown(f"""
        <div style="background: rgba(56, 189, 248, 0.1); padding: 20px; border-radius: 10px; border-left: 4px solid #38bdf8;">
            {st.session_state['ai_analysis']}
        </div>
        """, unsafe_allow_html=True)

    with col_export:
        st.markdown('<div class="glass-card" style="text-align: center;">', unsafe_allow_html=True)
        st.markdown("#### Export Pack")
        
        # PPTX Generation Logic
        if st.session_state['ai_analysis'] != "Click 'Generate' to run the Commentary Engine.":
            pptx_file = create_presentation(metrics_for_ppt, st.session_state['ai_analysis'])
            
            st.download_button(
                label="📥 Download Slides (.pptx)",
                data=pptx_file,
                file_name="FinSight_Board_Pack_Jun24.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                type="primary"
            )
            st.caption("Compatible with Google Slides")
        else:
            st.warning("Run AI Analysis first to enable export.")
        st.markdown('</div>', unsafe_allow_html=True)

def data_engine_view():
    st.markdown("### ⚙️ Data Engine & Governance")
    c1, c2 = st.columns(2)
    with c1:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.markdown("#### 🔌 Pipeline Health")
        st.success("NetSuite Connector: ACTIVE (45ms)")
        st.success("Data Warehouse: SYNCED (09:41 AM)")
        st.info("Anomaly Detection: RUNNING")
        st.markdown('</div>', unsafe_allow_html=True)
    with c2:
        st.markdown('<div class="glass-card">', unsafe_allow_html=True)
        st.markdown("#### 🛡️ Version Control")
        st.markdown("**Current Version:** v3.4.1 (Prod)")
        st.markdown("**Last Audit:** User 'Admin' approved forecast adjust.")
        st.markdown('</div>', unsafe_allow_html=True)
        
    st.dataframe(ANNUAL_DATA.head(), use_container_width=True)

# --- MAIN LAYOUT ---
st.markdown("""
<div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 20px;">
    <div>
        <h1 style="margin:0; font-size: 2.2rem; background: -webkit-linear-gradient(#fff, #94a3b8); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">FinSight Core</h1>
        <p style="color: #64748b; margin:0;">Enterprise Edition • Live Connection</p>
    </div>
    <div style="text-align: right;">
        <span style="background: rgba(56, 189, 248, 0.1); color: #38bdf8; padding: 5px 10px; border-radius: 15px; font-size: 0.8rem; font-weight: 600;">Connected: NetSuite OneWorld</span>
    </div>
</div>
""", unsafe_allow_html=True)

tab1, tab2 = st.tabs(["Executive Dashboard", "Data Engine"])

with tab1:
    dashboard_view()

with tab2:
    data_engine_view()
